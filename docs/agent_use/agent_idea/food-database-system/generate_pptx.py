#!/usr/bin/env python3
"""Generate Food Database System PPTX - v2"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Colors
ORANGE  = RGBColor(0xFB, 0x92, 0x3C)
GREEN   = RGBColor(0x4A, 0xDE, 0x80)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
YELLOW  = RGBColor(0xFB, 0xBF, 0x24)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x5A, 0x64, 0x78)
DARK_BG = RGBColor(0x08, 0x08, 0x0D)
SURFACE = RGBColor(0x0E, 0x10, 0x19)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_tag(slide, text, left, top, color, text_color=WHITE):
    rect = add_rect(slide, left, top, len(text)*0.12+0.5, 0.38, fill_color=color)
    txBox = slide.shapes.add_textbox(Inches(left+0.08), Inches(top+0.07), Inches(len(text)*0.12+0.3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(8)
    run.font.color.rgb = text_color
    return rect

# ─────────────────────────────────────────────
# SLIDE 1 — SPLASH
# ─────────────────────────────────────────────
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)

# Radial gradient effect via overlay rects
add_rect(s1, 0, 0, 13.33, 7.5, fill_color=RGBColor(0x08, 0x08, 0x0D))

add_tag(s1, "情境簡報 v2", 0.5, 0.8, RGBColor(0x25, 0x12, 0x04), ORANGE)

add_textbox(s1, "冰箱", 0.5, 1.8, 6, 1.5,
            font_size=60, bold=True, color=ORANGE)
add_textbox(s1, "×", 4.5, 1.9, 1, 1,
            font_size=60, bold=True, color=WHITE)
add_textbox(s1, "餐桌", 5.3, 1.8, 6, 1.5,
            font_size=60, bold=True, color=GREEN)

add_textbox(s1, "食物資料庫系統", 0.5, 3.3, 12, 0.8,
            font_size=28, bold=True, color=WHITE)

add_textbox(s1, "三種使用模式：節源、買買買、AI對話。\n不管想吃什麼，系統都幫你算出缺什麼。",
            0.5, 4.1, 10, 1,
            font_size=16, color=MUTED)

tags = ["節源模式", "買買買模式", "AI對話模式", "Expand 推薦引擎"]
for i, t in enumerate(tags):
    add_rect(s1, 0.5+i*2.4, 5.4, 2.2, 0.42,
             fill_color=RGBColor(0x14, 0x14, 0x1E),
             line_color=RGBColor(0x30, 0x30, 0x40))
    add_textbox(s1, t, 0.6+i*2.4, 5.45, 2.1, 0.35,
                font_size=11, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 2 — THREE MODES
# ─────────────────────────────────────────────
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)

add_tag(s2, "使用情境", 0.4, 0.4, RGBColor(0x25, 0x12, 0x04), ORANGE)
add_textbox(s2, "三種模式，覆蓋所有進廚房的理由",
            0.4, 0.85, 12, 0.6, font_size=24, bold=True, color=WHITE)
add_textbox(s2, "節源、規劃、探索——各有不同的輸入與輸出邏輯",
            0.4, 1.45, 10, 0.4, font_size=13, color=MUTED)

# Mode 1: 節源
add_rect(s2, 0.4, 1.95, 12.5, 1.35,
         fill_color=SURFACE,
         line_color=RGBColor(0x25, 0x12, 0x04))
add_rect(s2, 0.4, 1.95, 0.35, 1.35, fill_color=RGBColor(0x25, 0x12, 0x04))
add_textbox(s2, "1", 0.45, 2.4, 0.3, 0.4, font_size=14, bold=True, color=ORANGE)
add_textbox(s2, "節源模式 — 冰箱剩什麼，就做什麼",
            0.85, 2.05, 8, 0.35, font_size=13, bold=True, color=ORANGE)
add_textbox(s2, '用戶勾選：紅蘿蔔、高麗菜、蛋 → 系統回應：紅蘿蔔炒蛋 ✓ 完全符合、菇菇炒高麗菜 ✓ 完全符合',
            0.85, 2.45, 11.5, 0.5, font_size=11, color=MUTED)

# Mode 2: 買買買
add_rect(s2, 0.4, 3.4, 12.5, 1.35,
         fill_color=SURFACE,
         line_color=RGBColor(0x0D, 0x26, 0x15))
add_rect(s2, 0.4, 3.4, 0.35, 1.35, fill_color=RGBColor(0x0D, 0x26, 0x15))
add_textbox(s2, "2", 0.45, 3.85, 0.3, 0.4, font_size=14, bold=True, color=GREEN)
add_textbox(s2, "買買買模式 — 這週想吃什麼，一次規劃",
            0.85, 3.5, 8, 0.35, font_size=13, bold=True, color=GREEN)
add_textbox(s2, '用戶勾選這週想吃：滷肉飯、咖哩飯、煎餃 → 系統產出購物清單（已扣除庫存）',
            0.85, 3.9, 11.5, 0.5, font_size=11, color=MUTED)

# Mode 3: AI對話
add_rect(s2, 0.4, 4.85, 12.5, 2.1,
         fill_color=SURFACE,
         line_color=RGBColor(0x06, 0x1A, 0x20))
add_rect(s2, 0.4, 4.85, 0.35, 2.1, fill_color=RGBColor(0x06, 0x1A, 0x20))
add_textbox(s2, "3", 0.45, 5.3, 0.3, 0.4, font_size=14, bold=True, color=CYAN)
add_textbox(s2, "AI對話模式 — 不知道吃什麼，讓 AI 幫你想",
            0.85, 4.95, 9, 0.35, font_size=13, bold=True, color=CYAN)
chat = (
    '用戶：「我不知道要吃什麼」\n'
    'AI：「那你冰箱有什麼？」\n'
    '用戶：「高麗菜、紅蘿蔔、蛋、火鍋料」\n'
    'AI：「不如今天吃火鍋搭配蛋花湯怎麼樣？」\n'
    '用戶：「但我想多買一點」\n'
    'AI：「來挑挑有什麼喜歡吃的吧」（系統根據庫存展開相關推薦）'
)
add_textbox(s2, chat, 0.85, 5.35, 11.8, 1.4, font_size=10, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 3 — 8 SUB-DATABASES
# ─────────────────────────────────────────────
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)

add_tag(s3, "資料庫架構", 0.4, 0.35, RGBColor(0x25, 0x12, 0x04), ORANGE)
add_textbox(s3, "8 個子資料庫，撐起整個系統",
            0.4, 0.8, 12, 0.6, font_size=24, bold=True, color=WHITE)
add_textbox(s3, "每個資料庫解決一個問題，分工明確、可以各自擴充",
            0.4, 1.38, 10, 0.35, font_size=12, color=MUTED)

databases = [
    ("食材庫",         "每種食材的名稱、營養素、GI值、發炎指標、保存方式、保存期限",               "核心顆粒度",   ORANGE),
    ("食材狀態庫",     "使用者冰箱的目前庫存：數量、購買日期、過期預估、存放位置（冷藏/冷凍）",   "冰箱掃描",     ORANGE),
    ("食譜庫",         "每道菜的完整配方：所需食材列表、步驟、整備/烹飪時間、困難度",              "推薦主體",     GREEN),
    ("Expand 關聯庫",  "每道菜的「相關菜」清單：使用相同核心食材的其他料理、用戶歷史延伸",         "Expand 引擎",  GREEN),
    ("食材價格庫",     "每種食材的市場零售價格（每日更新）、批发價、特價週期、替代食材建議",       "成本控制",     CYAN),
    ("餐點成本計算",   "每道菜的食材成本加總 + 人力水電估算 → 真實成本與自己煮省了多少",          "省錢感知",     CYAN),
    ("餐廳外食庫",     "各餐廳的菜單、價格、食材組成、地點、等候時間",                             "外食替代",     CYAN),
    ("使用者偏好學習", "使用者喜歡/不喜歡的食材、做過的菜、歷史評分、近期購買記錄（Local 儲存）", "個人化排序",   GREEN),
]

col_colors = [ORANGE, ORANGE, GREEN, GREEN, CYAN, CYAN, CYAN, GREEN]

for i, (name, content, use, col) in enumerate(databases):
    row = i
    left_col = 0.4
    mid_col   = 4.5
    top = 1.85 + row * 0.67

    add_rect(s3, left_col, top, 3.8, 0.6,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s3, name, left_col+0.1, top+0.1, 3.6, 0.4,
                font_size=12, bold=True, color=col)

    add_rect(s3, mid_col, top, 7.1, 0.6,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s3, content, mid_col+0.1, top+0.08, 6.9, 0.5,
                font_size=10, color=MUTED)

    add_rect(s3, 11.7, top, 1.3, 0.6,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s3, use, 11.8, top+0.12, 1.2, 0.4,
                font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4 — EXPAND ENGINE
# ─────────────────────────────────────────────
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)

add_tag(s4, "Expand 推薦引擎", 0.4, 0.4, RGBColor(0x0D, 0x26, 0x15), GREEN)
add_textbox(s4, "核心邏輯：食材 Expandability",
            0.4, 0.85, 12, 0.6, font_size=24, bold=True, color=WHITE)
add_textbox(s4, "不是「你有什麼食材能做什麼菜」，而是「你選了這道菜，系統幫你展開相關的所有可能」",
            0.4, 1.45, 12, 0.4, font_size=13, color=MUTED)

# Flow boxes
flow_items = [
    ("選食材 or 選菜",    ORANGE),
    ("Expand 相關食譜",   GREEN),
    ("庫存鉤稽",          CYAN),
    ("計算差異",          PURPLE),
    ("產出購物清單",      GREEN),
]
for i, (txt, col) in enumerate(flow_items):
    x = 0.4 + i * 2.6
    add_rect(s4, x, 2.0, 2.2, 0.6, fill_color=RGBColor(0x0E, 0x10, 0x19), line_color=col)
    add_textbox(s4, txt, x+0.08, 2.1, 2.05, 0.45, font_size=10, bold=True, color=col)
    if i < len(flow_items)-1:
        add_textbox(s4, "→", x+2.2, 2.1, 0.4, 0.4, font_size=16, color=MUTED)

# Expand examples
examples = [
    ("用戶勾選\n紅蘿蔔、高麗菜",  "相關食譜：\n紅蘿蔔炒蛋、滷肉飯\n菇菇炒高麗菜、炸豬排",    ORANGE, GREEN),
    ("用戶選了\n滷肉飯",         "相關推薦：\n紅燒肉、咖哩飯\n（用相同蛋白質来源）",     ORANGE, PURPLE),
    ("用戶選了\n炸豬排",         "高麗菜絲 ✓ 已有\n豬排 ✗ 需購買\n麵包粉 ✗ 需購買",         ORANGE, CYAN),
]
for i, (core, result, c_col, r_col) in enumerate(examples):
    x = 0.4 + i * 4.3
    add_rect(s4, x, 2.9, 3.8, 1.35,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_rect(s4, x, 2.9, 1.3, 1.35, fill_color=RGBColor(0x25, 0x12, 0x04))
    add_textbox(s4, core, x+0.06, 3.05, 1.2, 1.1, font_size=10, bold=True, color=c_col, align=PP_ALIGN.CENTER)
    add_textbox(s4, result, x+1.35, 2.98, 2.4, 1.2, font_size=10, color=r_col)

add_textbox(s4, "Expand = 根據已選食材或已選菜，找出所有「使用相同核心食材」的其他料理，越多相關 = 越高 Expandability",
            0.4, 4.45, 12.5, 0.4, font_size=11, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 5 — THREE DIMENSIONS
# ─────────────────────────────────────────────
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)

add_tag(s5, "推薦維度", 0.4, 0.4, RGBColor(0x06, 0x1A, 0x20), CYAN)
add_textbox(s5, "三個維度同時評估，決定推薦排序",
            0.4, 0.85, 12, 0.6, font_size=24, bold=True, color=WHITE)
add_textbox(s5, "系統同時計算這三個分數，加權後輸出最終排序",
            0.4, 1.45, 10, 0.35, font_size=13, color=MUTED)

dims = [
    ("📦", "庫存覆蓋度",   "這道菜有多少%食材家裡已經有？100% = 完全符合，不需要買",        ORANGE),
    ("🔗", "Expandability","選這道之後，相關延伸料理有多少%家裡也有？越高越容易持續擴展",   GREEN),
    ("⏰", "時間窗口",     "整備時間 + 烹飪時間，符合使用者現在的時間嗎？",                  CYAN),
]
bar_widths = [4.9, 4.2, 3.5]
bar_weights = ["70%", "60%", "50%"]
for i, (icon, title, desc, col) in enumerate(dims):
    x = 0.4 + i * 4.3
    add_rect(s5, x, 1.95, 3.9, 2.5,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s5, icon, x+1.5, 2.05, 1, 0.6, font_size=28, align=PP_ALIGN.CENTER)
    add_textbox(s5, title, x+0.1, 2.65, 3.7, 0.4, font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s5, desc, x+0.15, 3.1, 3.65, 0.9, font_size=10, color=MUTED)

# Bar chart
add_textbox(s5, "維度加權（系統預設，可由使用者調整）", 0.4, 4.6, 8, 0.35, font_size=12, bold=True, color=WHITE)
bar_data = [("庫存覆蓋度", 4.9, ORANGE), ("Expandability", 4.2, GREEN), ("時間窗口", 3.5, CYAN)]
weights = ["70%", "60%", "50%"]
for i, (label, bw, col) in enumerate(bar_data):
    y = 5.05 + i * 0.52
    add_textbox(s5, label, 0.4, y, 2.5, 0.4, font_size=11, bold=True, color=col)
    add_rect(s5, 2.9, y+0.05, bw, 0.3, fill_color=col)
    add_textbox(s5, weights[i], 2.9+bw+0.1, y, 0.6, 0.4, font_size=11, color=MUTED)

add_textbox(s5, "使用者在「省錢模式」下可調高庫存覆蓋度權重，在「探索模式」下可調高 Expandability",
            0.4, 6.7, 12, 0.35, font_size=10, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 6 — ROADMAP
# ─────────────────────────────────────────────
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)

add_tag(s6, "實作方向", 0.4, 0.4, RGBColor(0x0D, 0x26, 0x15), GREEN)
add_textbox(s6, "三階段，循序漸進",
            0.4, 0.85, 12, 0.6, font_size=24, bold=True, color=WHITE)
add_textbox(s6, "不求一次到位，先做出能展示的 MVP",
            0.4, 1.45, 10, 0.35, font_size=13, color=MUTED)

phases = [
    ("①", "節源模式 MVP",       "食材庫(20種) + 食譜庫(30道)\n+ 庫存勾選 + 基礎推薦",              "JSON資料庫 + 可操作網頁",  ORANGE),
    ("②", "Expand + 買買買",    "Expand關聯庫 + 多菜規劃\n+ 購物清單生成 + 偏好Local學習",        "完整推薦邏輯 + 成本顯示",  GREEN),
    ("③", "AI對話 + 精誠串接",  "自然語言對話介面\n+ 精誠AI模型接入",                               "可展示APP原型",           CYAN),
]
for i, (num, title, desc, deliver, col) in enumerate(phases):
    x = 0.4 + i * 4.3
    add_rect(s6, x, 1.95, 3.9, 3.0,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s6, num, x+1.5, 2.1, 1, 0.7, font_size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s6, title, x+0.1, 2.85, 3.7, 0.4, font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s6, desc, x+0.15, 3.3, 3.65, 0.9, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s6, x+0.3, 4.25, 3.3, 0.01, fill_color=RGBColor(0x25, 0x25, 0x35))
    add_textbox(s6, deliver, x+0.15, 4.35, 3.65, 0.5, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_textbox(s6, "外部資源", 0.4, 5.15, 3, 0.35, font_size=11, bold=True, color=WHITE)
ext = [("🌐 Open Food Facts", "全球食材營養資料庫API"), ("🗺️ Google Maps", "餐廳地點、評價、外送"), ("🤖 精誠 AI 模型", "吃結構化資料，輸出自然語言")]
for i, (name, desc) in enumerate(ext):
    x = 0.4 + i * 4.3
    add_rect(s6, x, 5.55, 4.0, 1.4,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s6, name, x+0.15, 5.65, 3.7, 0.4, font_size=11, bold=True, color=GREEN)
    add_textbox(s6, desc, x+0.15, 6.1, 3.7, 0.7, font_size=10, color=MUTED)

# ─────────────────────────────────────────────
# SLIDE 7 — CLOSING
# ─────────────────────────────────────────────
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)

add_textbox(s7, "打開冰箱，", 0.5, 1.2, 12, 0.9, font_size=36, bold=True, color=WHITE)
add_textbox(s7, "就有答案", 0.5, 2.0, 12, 0.9, font_size=60, bold=True, color=ORANGE)

add_textbox(s7, "三種模式 + Expand引擎 + 庫存鉤稽，\n把「今天吃什麼」的問題，變成「還缺什麼」的清單。",
            0.5, 3.2, 10, 0.8, font_size=16, color=MUTED)

cards = [
    ("📐", "Schema 先行的價值",    "先定義資料長什麼樣子，Expand關聯庫做起來才有根據",    ORANGE),
    ("🚀", "MVP 先行",             "20種食材 + 30道食譜，就是第一個能展示的版本",          GREEN),
    ("🔗", "Local + AI 混合",     "偏好資料存本地，AI處理對話，兼顧隱私與智慧",          CYAN),
]
for i, (icon, title, desc, col) in enumerate(cards):
    x = 0.4 + i * 4.3
    add_rect(s7, x, 4.2, 4.0, 1.7,
             fill_color=SURFACE, line_color=RGBColor(0x20, 0x20, 0x30))
    add_textbox(s7, icon, x+1.5, 4.3, 1, 0.5, font_size=22, align=PP_ALIGN.CENTER)
    add_textbox(s7, title, x+0.1, 4.85, 3.8, 0.4, font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s7, desc, x+0.15, 5.3, 3.7, 0.55, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s7, 0.4, 6.1, 12.5, 0.9,
         fill_color=RGBColor(0x14, 0x10, 0x0A),
         line_color=ORANGE)
add_textbox(s7, "下一步：定義 Expand 關聯 Schema", 0.6, 6.18, 8, 0.35, font_size=13, bold=True, color=ORANGE)
add_textbox(s7, "先把「哪些菜跟哪些菜使用相同核心食材」定義清楚，第一筆資料就能立刻開始建",
            0.6, 6.55, 11.5, 0.35, font_size=10, color=MUTED)

# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "食物資料庫系統_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
